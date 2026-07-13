from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("specification-internet-deck.pptx")
PROTOTYPE = Path(__file__).with_name("prototype-auth.png")
SIRAJ = Path(__file__).with_name("siraj.jpg")
THANIS = Path(__file__).with_name("thanis.jpg")
SHERWIN = Path(__file__).with_name("sherwin.jpg")
FORM_PROCESS = Path(__file__).with_name("form-process-frustration.jpeg")
FORM_SEARCH = Path(__file__).with_name("form-supplier-search.jpeg")
FORM_AI_VALUE = Path(__file__).with_name("form-ai-value.jpeg")
FORM_MAKER = Path(__file__).with_name("form-maker-difficulty.jpeg")

INK = "0F1F1C"
DEEP = "12231F"
PAPER = "FBFCF8"
PAPER_2 = "EEF6F1"
LINE = "D9E7E0"
TEAL = "0F9F8F"
MINT = "72E0C9"
CORAL = "FF7A59"
AMBER = "F4C95D"
BLUE = "5F7DF2"
MUTED = "66756F"
WHITE = "FFFFFF"


def rgb(hex_color):
    return RGBColor.from_string(hex_color)


prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
BLANK = prs.slide_layouts[6]


def add_bg(slide, dark=False):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(DEEP if dark else PAPER)

    if dark:
        colors = [(MINT, 86), (BLUE, 92), (AMBER, 90)]
    else:
        colors = [(MINT, 78), (BLUE, 92), (AMBER, 88)]

    band = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(-1.5), Inches(1.65), Inches(10.8), Inches(0.72))
    band.rotation = -14
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(colors[1][0])
    band.fill.transparency = colors[1][1]
    band.line.fill.background()

    band2 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(9.3), Inches(6.25), Inches(8.7), Inches(0.78))
    band2.rotation = 16
    band2.fill.solid()
    band2.fill.fore_color.rgb = rgb(colors[2][0])
    band2.fill.transparency = colors[2][1]
    band2.line.fill.background()

    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.3), Inches(-0.9), Inches(5.8), Inches(4.2))
    glow.fill.solid()
    glow.fill.fore_color.rgb = rgb(colors[0][0])
    glow.fill.transparency = 82
    glow.line.fill.background()


def textbox(slide, text, x, y, w, h, size=18, bold=False, color=INK, font="Aptos", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def title(slide, text, y=1.55, w=9.1, dark=False, size=42):
    return textbox(slide, text, 0.78, y, w, 1.3, size=size, bold=True, color=WHITE if dark else INK, font="Georgia")


def eyebrow(slide, text, y=1.27, dark=False):
    return textbox(slide, text.upper(), 0.78, y, 8.2, 0.28, size=13, bold=True, color=MINT if dark else TEAL)


def brand(slide, n, dark=False):
    mark = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(0.49), Inches(0.43), Inches(0.43))
    mark.fill.solid()
    mark.fill.fore_color.rgb = rgb(BLUE)
    mark.line.fill.background()

    spine = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.88), Inches(0.58), Inches(0.07), Inches(0.25))
    spine.fill.solid()
    spine.fill.fore_color.rgb = rgb(WHITE)
    spine.line.fill.background()
    page = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.98), Inches(0.58), Inches(0.16), Inches(0.25))
    page.fill.solid()
    page.fill.fore_color.rgb = rgb(WHITE)
    page.line.fill.background()
    textbox(slide, "Briefly", 1.36, 0.58, 1.2, 0.24, size=15, bold=True, color=MINT if dark else TEAL)
    textbox(slide, f"{n:02d} / 18", 14.5, 0.63, 0.86, 0.22, size=11, bold=True, color="CFE9E1" if dark else MUTED, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title_text, bullets=None, dark=False, body=None, check=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(WHITE if not dark else "345049")
    shape.fill.transparency = 5 if not dark else 12
    shape.line.color.rgb = rgb(LINE if not dark else "506961")
    shape.line.width = Pt(0.8)
    textbox(slide, title_text, x + 0.18, y + 0.18, w - 0.35, 0.32, size=15, bold=True, color=WHITE if dark else INK)
    if body:
        textbox(slide, body, x + 0.18, y + 0.68, w - 0.35, h - 0.78, size=15, bold=True, color="DCECE7" if dark else INK)
    if bullets:
        bullet_list(slide, bullets, x + 0.2, y + 0.62, w - 0.4, h - 0.72, dark=dark, check=check)
    return shape


def bullet_list(slide, items, x, y, w, h, dark=False, check=False, size=11.5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("✓  " if check else "●  ") + item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = rgb("DCECE7" if dark else INK)
        p.space_after = Pt(5)
        if not check:
            p.runs[0].font.color.rgb = rgb("DCECE7" if dark else INK)
    return box


def quote(slide, text, x, y, w, h, dark=False, size=13, attribution=None):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.04), Inches(h))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(MINT if dark else TEAL)
    line.line.fill.background()
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.04), Inches(y), Inches(w - 0.04), Inches(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb("2A443E" if dark else WHITE)
    bg.fill.transparency = 5 if not dark else 18
    bg.line.fill.background()
    textbox(slide, text, x + 0.25, y + 0.2, w - 0.45, h - 0.22, size=size, bold=True, color=WHITE if dark else INK)
    if attribution:
        textbox(slide, attribution, x + 0.25, y + h - 0.35, w - 0.45, 0.2, size=11, bold=True, color="CFE9E1" if dark else MUTED)


def prototype(slide, x, y, w, h):
    outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    outer.fill.solid()
    outer.fill.fore_color.rgb = rgb(DEEP)
    outer.line.fill.background()
    for i, c in enumerate([CORAL, AMBER, MINT]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.25 + i * 0.17), Inches(y + 0.18), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(c)
        dot.line.fill.background()
    if PROTOTYPE.exists():
        slide.shapes.add_picture(str(PROTOTYPE), Inches(x + 0.18), Inches(y + 0.48), Inches(w - 0.36), Inches(h - 0.62))


def add_cropped_picture(slide, image_path, x, y, w, h, focus_x=0.5, focus_y=0.5):
    if not image_path.exists():
        return None

    from PIL import Image

    with Image.open(image_path) as im:
        image_w, image_h = im.size

    target_ratio = w / h
    image_ratio = image_w / image_h
    pic = slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), Inches(w), Inches(h))

    if image_ratio > target_ratio:
        crop = 1 - (target_ratio / image_ratio)
        pic.crop_left = crop * focus_x
        pic.crop_right = crop * (1 - focus_x)
    else:
        crop = 1 - (image_ratio / target_ratio)
        pic.crop_top = crop * focus_y
        pic.crop_bottom = crop * (1 - focus_y)

    return pic


def member_card(slide, x, y, w, h, image_path, name, nickname, role, note, focus_y=0.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(WHITE)
    shape.fill.transparency = 5
    shape.line.color.rgb = rgb(LINE)
    shape.line.width = Pt(0.8)

    add_cropped_picture(slide, image_path, x + 0.18, y + 0.18, w - 0.36, 2.95, focus_y=focus_y)

    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.26), Inches(y + 3.34), Inches(0.92), Inches(0.34))
    pill.fill.solid()
    pill.fill.fore_color.rgb = rgb(DEEP)
    pill.line.fill.background()
    textbox(slide, role, x + 0.43, y + 3.43, 0.5, 0.16, size=10, bold=True, color=MINT, align=PP_ALIGN.CENTER)

    textbox(slide, name, x + 0.26, y + 3.82, w - 0.52, 0.5, size=14.5, bold=True, color=INK)
    textbox(slide, nickname, x + 0.26, y + 4.36, w - 0.52, 0.18, size=10.5, bold=True, color=TEAL)
    textbox(slide, note, x + 0.26, y + 4.72, w - 0.52, 0.55, size=9.5, bold=True, color=MUTED)


def add_contained_picture(slide, image_path, x, y, w, h):
    if not image_path.exists():
        return None

    from PIL import Image

    with Image.open(image_path) as im:
        image_w, image_h = im.size

    image_ratio = image_w / image_h
    box_ratio = w / h
    if image_ratio > box_ratio:
        pic_w = w
        pic_h = w / image_ratio
    else:
        pic_h = h
        pic_w = h * image_ratio

    pic_x = x + (w - pic_w) / 2
    pic_y = y + (h - pic_h) / 2
    return slide.shapes.add_picture(str(image_path), Inches(pic_x), Inches(pic_y), Inches(pic_w), Inches(pic_h))


def chart_card(slide, x, y, w, h, heading, image_path):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(WHITE)
    shape.fill.transparency = 4
    shape.line.color.rgb = rgb(LINE)
    shape.line.width = Pt(0.8)
    textbox(slide, heading, x + 0.18, y + 0.18, w - 0.36, 0.25, size=13, bold=True, color=INK)
    add_contained_picture(slide, image_path, x + 0.18, y + 0.55, w - 0.36, h - 0.72)


def stat_card(slide, x, y, w, h, stat, heading, note):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(WHITE)
    shape.fill.transparency = 5
    shape.line.color.rgb = rgb(LINE)
    shape.line.width = Pt(0.8)
    textbox(slide, stat, x + 0.18, y + 0.16, w - 0.36, 0.28, size=18, bold=True, color=TEAL)
    textbox(slide, heading, x + 0.18, y + 0.55, w - 0.36, 0.34, size=13, bold=True, color=INK)
    textbox(slide, note, x + 0.18, y + 1.0, w - 0.36, h - 1.08, size=9.5, bold=True, color=MUTED)


def flow_nodes(slide, nodes, x, y, dark=False, horizontal=True):
    step_x = x
    step_y = y
    for i, node in enumerate(nodes):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(step_x), Inches(step_y), Inches(1.75), Inches(0.38))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb("405A53" if dark else WHITE)
        shape.fill.transparency = 8
        shape.line.color.rgb = rgb("637D75" if dark else LINE)
        shape.line.width = Pt(0.7)
        textbox(slide, node, step_x + 0.14, step_y + 0.1, 1.55, 0.22, size=10.5, bold=True, color=WHITE if dark else INK)
        if i < len(nodes) - 1:
            if horizontal:
                textbox(slide, "→", step_x + 1.95, step_y + 0.08, 0.22, 0.2, size=10, bold=True, color=CORAL)
                step_x += 2.25
            else:
                textbox(slide, "↓", step_x + 0.08, step_y + 0.48, 0.22, 0.2, size=10, bold=True, color=CORAL)
                step_y += 0.85


def table(slide, rows, x, y, w, row_h=0.44, first_w=4.0):
    total_h = row_h * len(rows)
    outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(total_h))
    outer.fill.solid()
    outer.fill.fore_color.rgb = rgb(WHITE)
    outer.line.color.rgb = rgb(LINE)
    outer.line.width = Pt(0.8)
    col_w = [first_w, (w - first_w) / 2, (w - first_w) / 2]
    for r, row in enumerate(rows):
        yy = y + r * row_h
        for c, cell in enumerate(row):
            xx = x + sum(col_w[:c])
            fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(xx), Inches(yy), Inches(col_w[c]), Inches(row_h))
            fill.fill.solid()
            fill.fill.fore_color.rgb = rgb(DEEP if r == 0 else WHITE)
            fill.line.color.rgb = rgb(LINE)
            fill.line.width = Pt(0.4)
            textbox(slide, cell, xx + 0.14, yy + 0.13, col_w[c] - 0.25, row_h - 0.16, size=10, bold=True, color=WHITE if r == 0 else INK)


def make_slide(n, dark=False):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, dark)
    brand(slide, n, dark)
    return slide


# 01
s = make_slide(1)
eyebrow(s, "The Specification Internet", 3.2)
title(s, "Commerce should\nbegin with human\nintent.", y=3.55, w=6.7, size=48)
textbox(s, "Not product catalogues. Briefly transforms what people want into\nstandardized manufacturing specifications.", 1.0, 5.72, 6.6, 0.62, size=15, bold=True)
quote(s, "Transforming human intent into standardized manufacturing\nspecifications.", 1.0, 6.48, 6.9, 0.78, size=13.5)
prototype(s, 8.2, 3.0, 7.2, 4.6)

# 02
s = make_slide(2)
eyebrow(s, "The hidden problem")
title(s, "Commerce begins with products. People do not.", w=9.5, size=36)
card(s, 0.78, 2.75, 6.9, 2.05, "Customers", ["Cannot find exactly what they need", "Compromise on design or functionality", "Contact multiple manufacturers", "Repeatedly explain requirements", "Wait days for quotations"])
card(s, 8.1, 2.75, 6.9, 2.05, "Manufacturers", ["Receive vague enquiries", "Spend hours clarifying requirements", "Educate customers on materials and constraints", "Prepare quotations that never convert"])
quote(s, "The real bottleneck is not manufacturing. It is translating human intent into something manufacturers can actually build.", 0.9, 5.35, 8.8, 0.78)

# 03
s = make_slide(3)
eyebrow(s, "The pivot")
title(s, "We thought we were solving the wrong problem.", w=9.5, size=36)
card(s, 0.78, 2.75, 6.9, 2.05, "Initial assumption")
flow_nodes(s, ["Find manufacturers", "Marketplace MVP"], 1.05, 3.38)
flow_nodes(s, ["5 interviews", "Real bottleneck"], 1.05, 4.05)
card(s, 8.1, 2.75, 6.9, 2.05, "What manufacturers told us", ["Customers do not know what manufacturers need", "Clarification wastes time", "Budgets and materials are misunderstood", "Standardized specs reduce quotation effort"], check=True)
quote(s, '"It turns an ambiguous design conversation into a deterministic procurement transaction."', 0.9, 5.35, 8.8, 0.9, attribution="Teck Soon Furniture")

# 04
s = make_slide(4, True)
eyebrow(s, "Introducing the Specification Internet", dark=True)
title(s, "We remove ambiguity before manufacturing begins.", w=9.5, dark=True, size=36)
card(s, 0.78, 2.75, 6.9, 2.6, "Instead of", dark=True)
flow_nodes(s, ["Customer", "Search products", "Compromise", "Purchase"], 1.08, 3.35, dark=True, horizontal=False)
card(s, 8.1, 2.75, 6.9, 2.6, "We propose", dark=True)
flow_nodes(s, ["Customer intent", "AI specification copilot", "Specification packet", "Multiple manufacturers"], 8.4, 3.35, dark=True, horizontal=False)
textbox(s, "We do not replace manufacturers. We help customers arrive with quote-ready requirements.", 0.9, 6.0, 9.0, 0.35, size=16, bold=True, color="DCECE7")

# 05
s = make_slide(5)
eyebrow(s, "The AI Specification Copilot")
title(s, "From a sentence to a manufacturable\npacket.", w=9.5, size=34)
card(s, 0.78, 3.8, 6.9, 3.1, "Customer intent")
textbox(s, '"I need a walnut study desk under my staircase."', 1.05, 4.37, 6.2, 0.28, size=11, bold=True)
card(s, 1.05, 5.0, 2.85, 1.6, "AI clarifies", ["Dimensions", "Budget", "Materials", "Doorway access", "Cable management", "Child-safe needs"], check=True)
card(s, 4.1, 5.0, 2.85, 1.6, "AI generates", ["Material recommendations", "Cost estimate", "Lead time estimate", "Manufacturability checks", "Simple 3D preview", "Specification packet"])
card(s, 8.0, 4.55, 7.1, 1.75, "Funding step change")
textbox(s, "Current MVP", 8.25, 5.12, 2.1, 0.25, size=12, bold=True)
bullet_list(s, ["Requirement collection", "Manufacturer discovery", "Quotation workflow"], 8.25, 5.45, 2.8, 0.7, check=True, size=9.5)
textbox(s, "Next phase", 11.85, 5.12, 2.1, 0.25, size=12, bold=True)
bullet_list(s, ["AI clarification engine", "Specification packet", "Material recommendation", "Cost estimation", "3D preview"], 11.85, 5.45, 2.8, 0.95, size=9.5)

# 06
s = make_slide(6)
eyebrow(s, "Our MVP")
title(s, "Build. Validate. Improve.", w=9.2, size=38)
card(s, 0.78, 4.1, 3.55, 1.85, "Current MVP", ["Browse custom furniture", "Submit requirements", "Receive quotations", "Compare manufacturers", "Place order"], check=True)
card(s, 4.55, 4.1, 3.55, 1.85, "Why we started here", ["Validate the workflow first", "Learn from real users", "Avoid unnecessary features", "Prove demand before scaling"])
quote(s, "Validate first. Build second. Scale third.", 0.9, 6.18, 7.2, 0.5, size=12)
prototype(s, 8.35, 3.0, 7.0, 4.8)

# 07
s = make_slide(7)
eyebrow(s, "Built. Tested. Pivoted.")
title(s, "We challenged our assumptions.", w=9.2, size=36)
flow_nodes(s, ["Problem", "Marketplace MVP", "Survey", "Manufacturer interviews", "Specification Internet"], 0.78, 2.55)
table(s, [
    ["User feedback", "Product improvement", "Strategic result"],
    ["Customers struggle to explain requirements", "AI clarification workflow", "Less vague demand"],
    ["Manufacturers ask repetitive questions", "Specification packet", "Faster quotation"],
    ["Specifications are inconsistent", "Universal structured format", "Comparable bids"],
    ["MVP too broad", "Focus on custom wooden furniture", "Sharper validation"],
], 0.78, 3.35, 14.3, row_h=0.43, first_w=4.0)

# 08
s = make_slide(8)
eyebrow(s, "Market validation")
title(s, "Validated by customers and\nmanufacturers.", w=8.0, size=38)
card(s, 0.78, 3.1, 6.9, 1.85, "Customer validation", ["XX respondents", "XX% could not find exactly what they wanted", "XX% compromised or gave up", "XX% found the AI workflow valuable", "XX waitlist signups"])
card(s, 8.1, 3.1, 6.9, 1.85, "Manufacturer validation", ["5/5 reported missing specifications", "5/5 spend significant time clarifying", "5/5 said standardized specs reduce effort", "4/5 interested in piloting"], check=True)
quote(s, '"If your AI can tell customers their budget does not match their material choice before they talk to me,\nyou save me 80% of my nightly headaches."', 0.9, 5.25, 8.6, 0.98, attribution="WD Custom Woodcraft", size=11.5)

# 09
s = make_slide(9)
eyebrow(s, "Why existing solutions fall short")
title(s, "Existing platforms help people find\nproducts. We help people define what\nthey need.", w=9.7, size=34)
table(s, [
    ["", "Existing solutions", "The Specification Internet"],
    ["Starting point", "Products", "Human intent"],
    ["Workflow", "Search and compromise", "AI clarification"],
    ["Specifications", "Manual", "Standardized"],
    ["Quotations", "One manufacturer", "Competitive bidding"],
    ["Innovation", "Marketplace", "Specification engine"],
], 0.78, 3.45, 14.3, row_h=0.39, first_w=4.0)

# 10
s = make_slide(10)
eyebrow(s, "Market opportunity")
title(s, "A focused wedge into custom\nmanufacturing.", w=8.0, size=37)
card(s, 0.78, 3.1, 4.7, 1.6, "Phase 1", body="Singapore custom wooden furniture")
card(s, 5.7, 3.1, 4.7, 1.6, "Phase 2", ["Wardrobes", "Kitchen cabinets", "Office furniture", "Retail displays"])
card(s, 10.62, 3.1, 4.7, 1.6, "Phase 3", ["Commercial fit-out", "Renovation", "Custom manufacturing"])
for x, label, body in [(0.78, "TAM", "Global custom manufacturing"), (3.25, "SAM", "Singapore custom furniture"), (5.72, "SOM", "500 users, 100 specs, 30 projects")]:
    k = slide = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.1), Inches(2.25), Inches(0.8))
    k.fill.solid()
    k.fill.fore_color.rgb = rgb(DEEP)
    k.line.fill.background()
    textbox(s, label, x + 0.16, 5.28, 1.1, 0.24, size=18, bold=True, color=MINT)
    textbox(s, body, x + 0.16, 5.62, 1.9, 0.18, size=9.5, bold=True, color=WHITE)
quote(s, "Long-term vision: a universal specification layer connecting customers and manufacturers\nworldwide.", 0.9, 6.12, 8.8, 0.72, size=11.5)

# 11
s = make_slide(11)
eyebrow(s, "Business model")
title(s, "Manufacturers pay when we create real\nbusiness.", w=10.4, size=38)
card(s, 0.78, 3.1, 6.9, 1.35, "Today: marketplace commission", body="Small commission on successfully completed projects.")
card(s, 8.1, 3.1, 6.9, 1.35, "Future: manufacturer SaaS", ["AI quotation assistant", "Specification management", "Workflow automation"])
card(s, 0.78, 4.75, 6.9, 1.35, "Enterprise API", body="Integrate the Specification Engine into manufacturing software.")
card(s, 8.1, 4.75, 6.9, 1.35, "Premium AI tools", ["Advanced design refinement", "Specification history", "Enhanced 3D previews"])

# 12
s = make_slide(12)
eyebrow(s, "Roadmap")
title(s, "From MVP to public beta.", w=8.0, size=38)
stages = [
    ("Already achieved", "Marketplace MVP,\nsurvey, interviews, pivot"),
    ("0-3 months", "AI clarification engine"),
    ("4-6 months", "Spec packet and material\nrecommendation"),
    ("7-9 months", "Pilot with 5-10\nmanufacturers"),
    ("10-12 months", "Public beta launch"),
]
for i, (head, body) in enumerate(stages):
    card(s, 0.78 + i * 2.9, 2.9, 2.65, 1.35, head, body=body)
for x, num, lab in [(0.78, "500", "users"), (2.55, "100", "specifications"), (4.32, "30", "completed projects"), (6.09, "5-10", "manufacturer partners")]:
    k = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.6), Inches(1.55), Inches(0.8))
    k.fill.solid()
    k.fill.fore_color.rgb = rgb(DEEP)
    k.line.fill.background()
    textbox(s, num, x + 0.15, 4.78, 0.9, 0.22, size=17, bold=True, color=MINT)
    textbox(s, lab, x + 0.15, 5.1, 1.25, 0.18, size=8.7, bold=True, color=WHITE)

# 13
s = make_slide(13)
eyebrow(s, "About us")
title(s, "Meet the team building Briefly.", w=9.2, size=38)
member_card(
    s,
    0.78,
    2.55,
    4.55,
    5.35,
    SIRAJ,
    "Mohamed Sirajudeen Bin Mohamed Rabik",
    "Siraj",
    "CEO",
    "Leads product direction, customer discovery, and the overall Briefly vision.",
    focus_y=0.56,
)
member_card(
    s,
    5.72,
    2.55,
    4.55,
    5.35,
    THANIS,
    "Thanis Senthilkumas",
    "Thanis",
    "CFO",
    "Owns finance, business modelling, and validation of the commercial strategy.",
    focus_y=0.42,
)
member_card(
    s,
    10.66,
    2.55,
    4.55,
    5.35,
    SHERWIN,
    "Sherwin Rufus",
    "Sherwin",
    "CTO",
    "Drives engineering, technical execution, and the AI specification roadmap.",
    focus_y=0.38,
)

# 14
s = make_slide(14)
eyebrow(s, "Why we are the right team")
title(s, "We started with a marketplace. Evidence\nchanged the product.", w=10.4, size=36)
card(s, 0.78, 3.1, 4.7, 1.35, "Technical", ["Full-stack development", "AI integration", "Product engineering"])
card(s, 5.7, 3.1, 4.7, 1.35, "Business", ["Customer discovery", "Validation", "Partnerships"])
card(s, 10.62, 3.1, 4.7, 1.35, "Design", ["UX/UI", "Prototyping", "User experience"])
quote(s, "Instead of defending our original idea, we changed it. Evidence of execution: working MVP, customer\nsurvey, manufacturer interviews, product pivot, AI roadmap, live demo.", 0.9, 4.85, 8.8, 0.72, size=10.8)

# 15
s = make_slide(15, True)
eyebrow(s, "The future of commerce", dark=True)
title(s, "Products built exactly as needed.", w=9.5, dark=True, size=38)
card(s, 0.78, 2.5, 6.9, 3.55, "Today", dark=True)
flow_nodes(s, ["Products", "Customers compromise"], 1.05, 3.15, dark=True, horizontal=False)
card(s, 8.1, 2.5, 6.9, 3.55, "Tomorrow", dark=True)
flow_nodes(s, ["Human intent", "AI spec copilot", "Standardized specification", "Manufacturers"], 8.38, 3.15, dark=True, horizontal=False)
quote(s, "The internet standardized how information moves. Manufacturing needs a standardized specification\nlayer so products begin with human intent, not product catalogues.", 0.9, 6.28, 8.8, 0.72, dark=True, size=10.8)
textbox(s, "We are not building a better marketplace. We are building the missing layer between\nimagination and manufacturing.", 0.78, 7.22, 8.7, 0.42, size=14, bold=True, color="DCECE7")


prs.save(OUT)
print(OUT)
